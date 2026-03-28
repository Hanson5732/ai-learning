import io
from fastapi import APIRouter, UploadFile, File, HTTPException
from pptx import Presentation
import PyPDF2  # 引入新安装的 PDF 解析库
from app.services.llm_service import extract_ppt_topics

router = APIRouter()

@router.post("/upload")
async def upload_and_parse_file(file: UploadFile = File(...)):
    # 1. 检查文件格式（把 .pdf 加上）
    filename = file.filename.lower()
    if not filename.endswith(('.ppt', '.pptx', '.pdf')):
        raise HTTPException(status_code=400, detail="目前只支持解析 .ppt, .pptx 或 .pdf 格式的文件哦！")
    
    try:
        content = await file.read()
        full_text = ""

        # 2. 分支处理：如果是 PDF
        if filename.endswith('.pdf'):
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text_runs = []
            for page in pdf_reader.pages:
                text = page.extract_text()
                if text:
                    text_runs.append(text.strip())
            full_text = "\n".join(text_runs)
            
        # 3. 分支处理：如果是 PPT
        else:
            prs = Presentation(io.BytesIO(content))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text.strip())
            full_text = "\n".join(text_runs)
        
        if not full_text.strip():
            raise HTTPException(status_code=400, detail="这个文件里好像全是图片，没有提取到文字内容哦。")
            
        # 4. 防爆 Token 截取
        safe_text = full_text[:4000]
        
        # 5. 召唤大模型干活
        llm_result = extract_ppt_topics(safe_text)
        
        return {
            "filename": file.filename,
            "topics": llm_result.get("topics", [])
        }
        
    except Exception as e:
        print(f"解析文件出错: {str(e)}")
        raise HTTPException(status_code=500, detail="文件解析失败，可能是文件损坏、加密或格式不兼容。")